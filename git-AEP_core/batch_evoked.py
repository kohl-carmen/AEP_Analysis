# hnn core
# https://hnnsolver.github.io/hnn-core/auto_examples/plot_simulate_evoked.html#sphx-glr-auto-examples-plot-simulate-evoked-py
# this does a variety of loops around the gaba parameters to explore their effects on the simulation waveform

#ipython
#-------------------------------------
#Init
#-------------------------------------
OG_HNN=False


## Original HNN
from neuron import h
h.nrn_load_dll('C:\\Users\\ckohl\\Miniconda3\\Lib\\site-packages\\hnn_core\\nrnmech.dll')
import os.path as op
if OG_HNN==True:
    #import hnn_core
    from hnn_core import simulate_dipole, Params, Network, read_params
else:
    #import hnn_core_ca
    from hnn_core_ca import simulate_dipole, Params, Network, read_params
import json
import numpy as np
import matplotlib.pyplot as plt
import sys        
#sys.path.append('C:\\Users\\ckohl\\Documents\\HNN core code\\')
sys.path.append('C:\\Users\\ckohl\\Documents\\git-AEP_core\\')
import my_hnn_core_functions as m
from pptx import Presentation #https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
from pptx.util import Inches

saving_dpl=False
my_param_out_dir='C:\\Users\\ckohl\\hnn_out\\'
dump_dir="C:\\Users\\ckohl\\Desktop\\Current\\HNN Core\\"     




#-------------------------------------
#PPT init
#-------------------------------------
left = 0
top=Inches(0.25)
height = Inches(7)
width=Inches(10)
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Set each evoked to zero"
subtitle.text = "batch_evoked.py"

#-------------------------------------
#Read base params and define params oi
#-------------------------------------
param_name='best_new_aep_r_contra'

dict_param=m.get_dict_param(my_param_out_dir,param_name)
base_params=Params(dict_param)
Inputs=[]
prox_count=0
dist_count=0
input_times=[]
for pi in base_params:
    if pi[0:8]=='t_evprox':
        prox_count=prox_count+1
        Inputs.append('Prox'+str(prox_count))
        input_times.append(base_params[pi])
    if pi[0:8]=='t_evdist':
        dist_count=dist_count+1
        Inputs.append('Dist'+str(dist_count))
        input_times.append(base_params[pi])
#sort inputs
sort_i=np.argsort(input_times)         
zipped_pairs = zip(sort_i,Inputs)       
Inputs = [x for _, x in sorted(zipped_pairs)]      

#evoked params
mylabel=['Spikes','t','sd','2Pyr_AMPA','2Pyr_NMDA','2Bask_AMPA','2Bas_NMDA','5Pyr_AMPA','5Pyr_NMDA','5Bask_AMPA','5Bask_NMDA']
actuallabel_pre=['numspikes','t','sigma_t','gbar','gbar','gbar','gbar','gbar','gbar','gbar','gbar']
actuallabel_post=['','','','_L2Pyr_ampa','_L2Pyr_nmda','_L2Basket_ampa','_L2Basket_nmda','_L5Pyr_ampa','_L5Pyr_nmda','_L5Basket_ampa','_L5Basket_nmda']
for input in range(0,len(Inputs)):
    these_params=[]
    #lets go through inputs and copy them
    if Inputs[input][0]=='P':
        it_through_these=mylabel
        actuallabel_mid='_evprox_'
        input_number=prox_count
    else:
        it_through_these=mylabel[0:9]
        actuallabel_mid='_evdist_'
        input_number=dist_count
        
    for input_p in range(0,len(it_through_these)):    
        these_params=base_params
            
        this_param=actuallabel_pre[input_p]+actuallabel_mid+Inputs[input][-1]+actuallabel_post[input_p]
           
        if mylabel[input_p]=='t':
            these_params[this_param]=these_params[this_param]+these_params[this_param]
        else:
            these_params[this_param]=0
            
        # run simulation

        net=[]
        dpls=[]
        net = Network(these_params)
        dpls = simulate_dipole(net, n_jobs=1, n_trials=10)
        # name='All - '+str(v*100)+'%'
        # m.core_output_basic(dpls,net,name,data,0,0,1,param_oi,base_params)   
        m.plot_hnn_core_output(dpls,net,Inputs[input]+'-'+this_param+' = 0',False,True,False,these_params,base_params)
        plt.pause(.1)
        plt.savefig(dump_dir+'\\temp.png')
        img_path = dump_dir+'\\temp.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
        plt.close()
        print(p)
        print(v)
        prs.save(dump_dir+'\\each_evoked_0.pptx')

        
