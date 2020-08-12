# hnn core
# https://hnnsolver.github.io/hnn-core/auto_examples/plot_simulate_evoked.html#sphx-glr-auto-examples-plot-simulate-evoked-py
# this does a variety of loops around the gaba parameters to explore their effects on the simulation waveform

#ipython
#-------------------------------------
#Init
#-------------------------------------
OG_HNN=False


## Original HNN
from neuron import h
h.nrn_load_dll('C:\\Users\\ckohl\\Miniconda3\\Lib\\site-packages\\hnn_core\\nrnmech.dll')
import os.path as op
if OG_HNN==True:
    #import hnn_core
    from hnn_core import simulate_dipole, Params, Network, read_params
else:
    #import hnn_core_ca
    from hnn_core_ca import simulate_dipole, Params, Network, read_params
import json
import numpy as np
import matplotlib.pyplot as plt
import sys        
#sys.path.append('C:\\Users\\ckohl\\Documents\\HNN core code\\')
sys.path.append('C:\\Users\\ckohl\\Documents\\git-AEP_core\\')
import my_hnn_core_functions as m
from pptx import Presentation #https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
from pptx.util import Inches

saving_dpl=False
my_param_out_dir='C:\\Users\\ckohl\\hnn_out\\'
dump_dir="C:\\Users\\ckohl\\Desktop\\Current\\HNN Core\\"     




#-------------------------------------
#PPT init
#-------------------------------------
left = 0
top=Inches(0.25)
height = Inches(7)
width=Inches(10)
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "New GABA test"
subtitle.text = "hnn_core_loop.py"


   
#-------------------------------------
#Read base params and define params oi
#-------------------------------------
param_name='best_new_aep_r_contra'

dict_param=m.get_dict_param(my_param_out_dir,param_name)
param_oi=['gbar_L2Basket_L2Pyr_gabaa','gbar_L2Basket_L2Pyr_gabab','gbar_L2Basket_L5Pyr','gbar_L5Basket_L5Pyr_gabaa','gbar_L5Basket_L5Pyr_gabab','gbar_L2Basket_L2Basket','gbar_L5Basket_L5Basket']
origin_layer=['L2', 'L2', 'L2', 'L5','L5','L2','L5']
connection=['I_E', 'I_E','I_E','I_E','I_E','I_I','I_I']
destination_layer=['L2','L2', 'L5', 'L5','L5','L2','L5']
destination_morph=['soma','soma','dendrite','soma','soma','singlecompartment','singlecompartment']
speed=['fast','slow','fast','fast','slow','fast','fast']        

Speed={'fast','slow'}
Connection={'I_E','I_I'}
Destination_Moprh={'soma','dendrite'}
Or_Layer={'L2','L5'}
Dest_Layer={'L2','L5'}

var_by=[0, .1, .2,.5, .8, 1]


#-------------------------------------------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------------------------------------
# Looping through param values
#-------------------------------------------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------------------------------------------

#-------------------------------------
#Vary all parameters
#-------------------------------------
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "TestEach new base: newtiina"
subtitle.text = "hnn_core_loop.py"
data='C:\\Users\\ckohl\\Documents\\SourcePlay\\NewTiina\\Adult_Tone-L Hemi-r.txt'
for v in var_by:
    base_params=[]
    these_params=[]
    dict_param=m.get_dict_param(my_param_out_dir,param_name)
    base_params = Params(dict_param)
    these_params=Params(dict_param)
    net=[]
    dpls=[]
    for p in range(0,len(param_oi)):
        these_params[param_oi[p]]=base_params[param_oi[p]]*(1-v)
    net = Network(these_params)
    dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
    # name='All - '+str(v*100)+'%'
    # m.core_output_basic(dpls,net,name,data,0,0,1,param_oi,base_params)   
    m.plot_hnn_core_output(dpls,net,'All - '+str(v*100)+'%',False,False,True,param_oi,these_params,base_params)
    plt.pause(.1)
    plt.savefig(dump_dir+'\\temp.png')
    img_path = dump_dir+'\\temp.png'
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
    plt.close()
    print(p)
    print(v)
prs.save(dump_dir+'\\newtiina_test_all1.pptx')

#-------------------------------------
#Vary each parameter separately
#-------------------------------------
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "TestEach new base: newtiina"
subtitle.text = "hnn_core_loop.py"
for p in param_oi:
    for v in var_by:
        base_params=[]
        these_params=[]
        dict_param=m.get_dict_param(my_param_out_dir,param_name)
        base_params = Params(dict_param)
        these_params=Params(dict_param)
        net=[]
        dpls=[]
        these_params[p]=base_params[p]*(1-v)
        net = Network(these_params)
        dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
        m.plot_hnn_core_output(dpls,net,p+' - '+str(v*100)+'%',False,False,True,param_oi,these_params,base_params)
        plt.pause(.1)
        plt.savefig(dump_dir+'\\temp.png')
        img_path = dump_dir+'\\temp.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
        plt.close()
        print(p)
        print(v)
prs.save(dump_dir+'\\newtiina_test_each1.pptx')




#-------------------------------------
#Vary all fast then all slow (all GABAA vs all GABAB)
#-------------------------------------
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "fastvsslow new base: newtiina"
subtitle.text = "hnn_core_loop.py"
for speed_i in Speed:
     for v in var_by:
        base_params=[]
        these_params=[]
        dict_param=m.get_dict_param(my_param_out_dir,param_name)
        base_params = Params(dict_param)
        these_params=Params(dict_param)
        net=[]
        dpls=[]
        for p in range(0,len(param_oi)):
            if speed[p]==speed_i:
                these_params[param_oi[p]]=base_params[param_oi[p]]*(1-v)
        net = Network(these_params)
        dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
        m.plot_hnn_core_output(dpls,net,'All '+speed_i+' - '+str(v*100)+'%',False,False,True,param_oi,these_params,base_params)
        plt.pause(.1)
        plt.savefig(dump_dir+'\\temp.png')
        img_path = dump_dir+'\\temp.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
        plt.close()
        print(p)
        print(v)
prs.save(dump_dir+'\\newtiina_fastvsslow1.pptx')


#-------------------------------------
#Vary all I->E connections vs all I->I connections
#-------------------------------------
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "IE/II new base: newtiina"
subtitle.text = "hnn_core_loop.py"
for conn_i in Connection:
     for v in var_by:
        base_params=[]
        these_params=[]
        dict_param=m.get_dict_param(my_param_out_dir,param_name)
        base_params = Params(dict_param)
        these_params=Params(dict_param)
        net=[]
        dpls=[]
        for p in range(0,len(param_oi)):
            if connection[p]==conn_i:
                these_params[param_oi[p]]=base_params[param_oi[p]]*(1-v)
        net = Network(these_params)
        dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
        m.plot_hnn_core_output(dpls,net,'All '+conn_i+' - '+str(v*100)+'%',False,False,True,param_oi,these_params,base_params)
        plt.pause(.1)
        plt.savefig(dump_dir+'\\temp.png')
        img_path = dump_dir+'\\temp.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
        plt.close()
        print(p)
        print(v)
prs.save(dump_dir+'\\newtiina_IE_II1.pptx')



#-------------------------------------
#Vary all soma connections vs all dendrite conenctions
#-------------------------------------
var_by=[0, .1, .2,.5, .8, 1]
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Morph new base:  newtiina"
subtitle.text = "hnn_core_loop.py"
for morph_i in Destination_Moprh:
     for v in var_by:
        base_params=[]
        these_params=[]
        dict_param=m.get_dict_param(my_param_out_dir,param_name)
        base_params = Params(dict_param)
        these_params=Params(dict_param)
        net=[]
        dpls=[]
        for p in range(0,len(param_oi)):
            if destination_morph[p]==morph_i:
                these_params[param_oi[p]]=base_params[param_oi[p]]*(1-v)
        net = Network(these_params)
        dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
        m.plot_hnn_core_output(dpls,net,'All '+morph_i+' - '+str(v*100)+'%',False,False,True,param_oi,these_params,base_params)
        plt.pause(.1)
        plt.savefig(dump_dir+'\\temp.png')
        img_path = dump_dir+'\\temp.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
        plt.close()
        print(p)
        print(v)
prs.save(dump_dir+'\\newtiina_morph1.pptx')



#-------------------------------------
#Vary all conenctions coming from L2 vs L5
#-------------------------------------
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Origin_Layer new  new base: newtiina"
subtitle.text = "hnn_core_loop.py"
for or_i in Or_Layer:
     for v in var_by:
        base_params=[]
        these_params=[]
        dict_param=m.get_dict_param(my_param_out_dir,param_name)
        base_params = Params(dict_param)
        these_params=Params(dict_param)
        net=[]
        dpls=[]
        for p in range(0,len(param_oi)):
            if origin_layer[p]==or_i:
                these_params[param_oi[p]]=base_params[param_oi[p]]*(1-v)
        net = Network(these_params)
        dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
        m.plot_hnn_core_output(dpls,net,'All '+or_i+' - '+str(v*100)+'%',False,False,True,param_oi,these_params,base_params)
        plt.pause(.1)
        plt.savefig(dump_dir+'\\temp.png')
        img_path = dump_dir+'\\temp.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
        plt.close()
        print(p)
        print(v)
prs.save(dump_dir+'\\newtiina_originlayer1.pptx')



#-------------------------------------
#Vary all conenctions going to L2 vs L5
#-------------------------------------
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Destination_Layer new  new base: newtiina"
subtitle.text = "hnn_core_loop.py"
for dest_i in Dest_Layer:
     for v in var_by:
        base_params=[]
        these_params=[]
        dict_param=m.get_dict_param(my_param_out_dir,param_name)
        base_params = Params(dict_param)
        these_params=Params(dict_param)
        net=[]
        dpls=[]
        for p in range(0,len(param_oi)):
            if destination_layer[p]==dest_i:
                these_params[param_oi[p]]=base_params[param_oi[p]]*(1-v)
        net = Network(these_params)
        dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
        m.plot_hnn_core_output(dpls,net,'All '+dest_i+' - '+str(v*100)+'%',False,False,True,param_oi,these_params,base_params)
        plt.pause(.1)
        plt.savefig(dump_dir+'\\temp.png')
        img_path = dump_dir+'\\temp.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
        plt.close()
        print(p)
        print(v)
prs.save(dump_dir+'\\newtiina_destinationlayer1.pptx')



#-------------------------------------
#Vary each parameter separately but increase (!)
#-------------------------------------
var_by=[0, 1, 2,5, 8, 10]
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Increase each  new base:newtiina"
subtitle.text = "hnn_core_loop.py"
for p in param_oi:
    for v in var_by:
        base_params=[]
        these_params=[]
        dict_param=m.get_dict_param(my_param_out_dir,param_name)
        base_params = Params(dict_param)
        these_params=Params(dict_param)
        net=[]
        dpls=[]
        these_params[p]=base_params[p]*(1+(v))
        net = Network(these_params)
        dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
        m.plot_hnn_core_output(dpls,net,p+' + '+str(v*100)+'%',False,False,True,param_oi,these_params,base_params)
        plt.pause(.1)
        plt.savefig(dump_dir+'\\temp.png')
        img_path = dump_dir+'\\temp.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
        plt.close()
        print(p)
        print(v)
prs.save(dump_dir+'\\newtiinatest_each_increasehigh1.pptx')

# #save
# if saving_dpl==True:
    # import pickle
    # f = open(dump_dir+param_name+'_dpl.txt', 'wb')
    # pickle.dump(dpls, f)
    # f.close()
    # f = open(dump_dir+param_name+'_net.txt', 'wb')
    # pickle.dump(net, f)
    # f.close()







# #-------------------------------------------------------------------------------------------------------------------------------
# #-------------------------------------------------------------------------------------------------------------------------------
# # Optimisation
# #-------------------------------------------------------------------------------------------------------------------------------
# #-------------------------------------------------------------------------------------------------------------------------------

# from scipy.optimize import minimize
# param_oi=['gbar_L2Basket_L2Pyr_gabaa','gbar_L2Basket_L2Pyr_gabab','gbar_L2Basket_L5Pyr','gbar_L5Basket_L5Pyr_gabaa','gbar_L5Basket_L5Pyr_gabab','gbar_L2Basket_L2Basket','gbar_L5Basket_L5Basket']
# my_param_out_dir='C:\\Users\\ckohl\\hnn_out\\'
# param_name='from_here_10_opt_adjusted'
# dict_param=m.get_dict_param(my_param_out_dir,param_name)
# base_params = Params(dict_param)
# x0=np.empty([1,len(param_oi)])
# for p in range(0,len(param_oi)):
    # x0[0,p]=base_params[param_oi[p]]   
# import fun_for_opt as f
# import time
# t = time.time()
# res = minimize(f.fun_for_opt, x0, method='nelder-mead',options={'maxiter': 50, 'disp': True})
# elapsed = time.time() - t

# #plot optimisation results:
# base_params = Params(dict_param)
# these_params=Params(dict_param)
# net=[]
# dpls=[]
# for p in range(0,len(param_oi)):
    # these_params[param_oi[p]]=res2['x'][p]
# net = Network(these_params)
# dpls = simulate_dipole(net, n_jobs=1, n_trials=3)
# m.plot_hnn_core_output(dpls,net,'Optimise',True,param_oi,these_params,base_params)     



# import pickle
# f = open(dump_dir+param_name+'_res.txt', 'wb')
# pickle.dump(res, f)
# f.close()

# f.close()




# pickle.read(f)













# import time
# t = time.time()
# net = Network(these_params,n_jobs=2)
# dpls = simulate_dipole(net, n_jobs=1, n_trials=2)
# elapsed = time.time() - t





 
