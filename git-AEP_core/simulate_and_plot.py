# hnn core
# for this to run, remember activate hnn or use Python Interpreter: Python 3.7.9 64 64-bit ('hnn': conda)

# this just simulates based on a given param file and plots output
# you can choose to put output plot in ppt and to save the dpl 

ppt=True
save_dpl=False
param_files=['best_new_aep_r_contra']
param_files=['best_aep_l_contra2']
data='C:\\Users\\ckohl\\Desktop\\Current\\AEP\\Data\\Adult_Tone-R Hemi-l.txt'

#-------------------------------------
#Init
#-------------------------------------
import os
os.chdir('C:\\Users\\ckohl\\hnn-core\\')
from neuron import h
h.nrn_load_dll('C:\\Users\\ckohl\\Miniconda3\\Lib\\site-packages\\hnn_core\\nrnmech.dll')
import hnn_core
from hnn_core import simulate_dipole, Params, Network, read_params
import numpy as np
import matplotlib.pyplot as plt
import sys        
sys.path.append('C:\\Users\\ckohl\\Desktop\\Current\\AEP\\git-AEP_Analysis\\git-AEP_core\\')
import my_hnn_core_functions as m
if ppt:
    from pptx import Presentation #https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
    from pptx.util import Inches
    #PPT init
    left = 0
    top=Inches(0.25)
    height = Inches(7)
    width=Inches(10)
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AEP"
    subtitle.text = "simulate_and_plot.py"


save_base_dir="C:\\Users\\ckohl\\Desktop\\Current\\AEP\\"     
my_param_out_dir=save_base_dir+'Params\\'#'C:\\Users\\ckohl\\hnn_out\\'
my_param_out_dir='C:\\Users\\ckohl\\hnn_out\\'

for p in param_files:
        dict_param=m.get_dict_param(my_param_out_dir,p)
        these_params=[]
        these_params=Params(dict_param)
        net=[]
        dpls=[]
        net = Network(these_params)
        dpls = simulate_dipole(net, n_trials=1)
        m.core_output_basic(dpls,net,p,data,False,False,False,False)
        if ppt:
            plt.pause(.1)
            plt.savefig(save_base_dir+'\\temp.png')
            img_path = save_base_dir+'\\temp.png'
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
            plt.close()

if ppt:
    prs.save(save_base_dir+'\\simulate_and_plot_output.pptx')




if save_dpl:
    import pickle
    f = open(save_base_dir+p+'_dpl.txt', 'wb')
    pickle.dump(dpls, f)
    f.close()
    f = open(save_base_dir+p+'_net.txt', 'wb')
    pickle.dump(net, f)
    f.close()




# paramoi=[]
# for pt in l:
#     if l[pt]!= r[pt]:
#         paramoi.append(pt)
# for pt in paramoi:
#     print(pt)
#     print('Difference:'+str(np.round(((l[pt]/r[pt])*100)-100))+'%')
#     print('')